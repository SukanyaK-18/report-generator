"""
Deployment Report Automation - Streamlit Web App
Generates a single-slide PowerPoint with a status table and clustered bar chart.
Run with: streamlit run report_generator/deployment_report_web.py
"""

import os
import io
import pandas as pd
import numpy as np
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from datetime import datetime

# --------------------------
# Page Config
# --------------------------
st.set_page_config(page_title="Deployment Report Generator", page_icon="📊", layout="wide")

# Colors matching the screenshot
HEADER_BLUE = RGBColor(0x44, 0x72, 0xC4)  # Blue Accent 1 (Office 2013-2022)
CHART_BLUE = RGBColor(0x44, 0x72, 0xC4)   # Blue bars
CHART_ORANGE = RGBColor(0xED, 0x7D, 0x31)  # Orange bars
CHART_COLORS = [
    CHART_BLUE, CHART_ORANGE,
    RGBColor(0xA9, 0xD1, 0x8E),
    RGBColor(0xFF, 0x6B, 0x6B),
    RGBColor(0x9B, 0x59, 0xB6),
    RGBColor(0x1A, 0xBC, 0x9C),
]


# --------------------------
# Chart Helper
# --------------------------
def set_category_axis_horizontal(chart):
    """Set category axis labels to horizontal (no rotation)."""
    catAx = chart._element.find(".//" + qn("c:catAx"))
    if catAx is not None:
        existing_txPr = catAx.find(qn("c:txPr"))
        if existing_txPr is not None:
            catAx.remove(existing_txPr)
        txPr = etree.SubElement(catAx, qn("c:txPr"))
        bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
        bodyPr.set("rot", "0")
        bodyPr.set("vert", "horz")
        etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", "900")


# --------------------------
# Build Single-Slide PPT
# --------------------------
def build_ppt(pivot, today_str):
    """
    Build a single-slide PowerPoint matching the exact format:
    - Header: "ATC-I POS Deployment/Remediation/TA status | {date}"
    - Bullet: "Deployments Weekly Status"
    - Table: Service | Start Date | End Date | Total Lanes | Success | Auto Completed | Offline | Remediated | Not Applicable | Unable to remediate | % of completion
    - Clustered bar chart: grouped by status category, one series per Service
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # ── Bullet title ─────────────────────────────────────
    title_tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.4), Inches(9.4), Inches(0.4))
    title_p = title_tb.text_frame.paragraphs[0]
    title_p.text = "•  Deployments  Weekly Status"
    title_p.font.size = Pt(14)
    title_p.font.bold = True

    # ── Define table columns ─────────────────────────────
    # Column order matching the screenshot exactly
    status_order = ["Success", "Auto Completed", "Offline", "Remediated", "Not Applicable", "Unable to remediate"]
    table_cols = ["Service", "Start Date", "End Date", "Total Lanes"] + status_order + ["% of completion"]

    # Map display names to actual pivot column names
    status_col_map = {
        "Success": "Completed",
        "Auto Completed": "Auto Completed",
        "Offline": "Offline",
        "Remediated": "Remediated",
        "Not Applicable": "Not Applicable",
        "Unable to remediate": "Unable to Fix",
    }

    n_rows = len(pivot) + 1  # +1 for header
    n_cols = len(table_cols)

    # ── Add table ────────────────────────────────────────
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.2), Inches(0.9),
        Inches(9.6), Inches(0.3 + len(pivot) * 0.35)
    )
    table = table_shape.table

    # Set column widths
    col_widths = [Inches(2.2), Inches(0.8), Inches(0.8), Inches(0.6),
                  Inches(0.7), Inches(0.8), Inches(0.6), Inches(0.8),
                  Inches(0.9), Inches(0.9), Inches(0.9)]
    for i, w in enumerate(col_widths[:n_cols]):
        table.columns[i].width = w

    # Header row
    for j, col_name in enumerate(table_cols):
        cell = table.cell(0, j)
        cell.text = col_name
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BLUE

    # Data rows
    for i, (_, row) in enumerate(pivot.iterrows()):
        for j, col_name in enumerate(table_cols):
            cell = table.cell(i + 1, j)
            if col_name == "Service":
                val = str(row.get("Deployment Type", ""))
            elif col_name == "Start Date":
                raw = row.get("Start Date", "")
                if pd.notna(raw):
                    if isinstance(raw, (pd.Timestamp, datetime)):
                        val = raw.strftime("%-m/%-d/%Y") if os.name != "nt" else raw.strftime("%#m/%#d/%Y")
                    else:
                        val = str(raw)
                else:
                    val = ""
            elif col_name == "End Date":
                raw = row.get("End Date", "")
                if pd.notna(raw):
                    if isinstance(raw, (pd.Timestamp, datetime)):
                        val = raw.strftime("%-m/%-d/%Y") if os.name != "nt" else raw.strftime("%#m/%#d/%Y")
                    else:
                        val = str(raw)
                else:
                    val = ""
            elif col_name == "Total Lanes":
                val = str(int(row.get("Total Lanes", 0)))
            elif col_name == "% of completion":
                total = int(row.get("Total Lanes", 0))
                offline = int(row.get(status_col_map.get("Offline", "Offline"), 0))
                success = int(row.get(status_col_map.get("Success", "Completed"), 0))
                auto = int(row.get(status_col_map.get("Auto Completed", "Auto Completed"), 0))
                remediated = int(row.get(status_col_map.get("Remediated", "Remediated"), 0))
                completed = success + auto + remediated
                denominator = total - offline
                pct = (completed / denominator * 100) if denominator > 0 else 0
                val = f"{pct:.2f}%"
            else:
                # Status columns — map display name to actual pivot column
                actual_col = status_col_map.get(col_name, col_name)
                val = str(int(row.get(actual_col, 0)))

            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9)
            p.alignment = 1  # Center
            # No background color on data rows
            cell.fill.background()

    # ── Chart title ────────────────────────────────────
    chart_title_top = Inches(0.9 + 0.35 * len(pivot) + 0.3)
    chart_title_tb = slide.shapes.add_textbox(Inches(0.3), chart_title_top, Inches(9.4), Inches(0.4))
    chart_title_p = chart_title_tb.text_frame.paragraphs[0]
    chart_title_p.text = "Deployments / Remediation-Weekly status"
    chart_title_p.font.size = Pt(12)
    chart_title_p.font.bold = True
    chart_title_p.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
    chart_title_p.alignment = PP_ALIGN.CENTER

    # ── Clustered bar chart ──────────────────────────────
    # X-axis: status categories (excluding Offline and Not Applicable from chart)
    chart_statuses = ["Success", "Auto Completed", "Offline", "Remediated", "Unable to remediate"]

    chart_data = CategoryChartData()
    chart_data.categories = chart_statuses

    for _, row in pivot.iterrows():
        service_name = str(row.get("Deployment Type", ""))
        values = [int(row.get(status_col_map.get(s, s), 0)) for s in chart_statuses]
        chart_data.add_series(service_name, values)

    chart_top = chart_title_top + Inches(0.4)
    chart_height = Inches(7.5) - chart_top - Inches(0.3)
    if chart_height < Inches(2.5):
        chart_height = Inches(2.5)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), chart_top,
        Inches(9.4), chart_height,
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    # Apply colors to each series (one per deployment)
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = CHART_COLORS[idx % len(CHART_COLORS)]

    # Data labels
    plot = chart.plots[0]
    plot.vary_by_categories = False  # Force legend to show series (deployment) names, not category names
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.show_category_name = False
    data_labels.number_format = '0'
    data_labels.font.size = Pt(9)

    # Style axes
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    set_category_axis_horizontal(chart)

    return prs


