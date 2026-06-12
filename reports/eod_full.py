"""
EOD Report Generator - Streamlit Web App
Generates a 2-slide PowerPoint for EOD Weekly Health Check from a single EOD Excel sheet.
"""

import io
import pandas as pd
import numpy as np
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from datetime import datetime

# Colors
HEADER_BLUE = RGBColor(0x44, 0x72, 0xC4)
CHART_BLUE = RGBColor(0x44, 0x72, 0xC4)
CHART_ORANGE = RGBColor(0xED, 0x7D, 0x31)
CHART_GRAY = RGBColor(0xA5, 0xA5, 0xA5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def get_device_type(hostname):
    """Derive device type (REG/SCO) from hostname."""
    hostname = str(hostname).upper()
    if "SCO" in hostname:
        return "SCO"
    elif "REG" in hostname:
        return "REG"
    return "Unknown"


def set_category_axis_horizontal(chart):
    catAx = chart._element.find(".//" + qn("c:catAx"))
    if catAx is not None:
        existing = catAx.find(qn("c:txPr"))
        if existing is not None:
            catAx.remove(existing)
        txPr = etree.SubElement(catAx, qn("c:txPr"))
        bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
        bodyPr.set("rot", "0")
        bodyPr.set("vert", "horz")
        etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", "900")


def compute_eod_data(df):
    """Compute all metrics from the raw EOD data."""
    # Identify columns
    status_col = None
    for c in df.columns:
        if "new/pre" in str(c).lower() or "pre-existing" in str(c).lower():
            status_col = c
            break

    category_col = "Category" if "Category" in df.columns else None
    issue_sub_col = "Issue Sub Category" if "Issue Sub Category" in df.columns else None

    # Summary stats
    total_lanes = len(df)
    offline_lanes = len(df[df[category_col].astype(str).str.strip().str.lower() == "offline"]) if category_col else 0
    remediated = len(df[df[status_col].astype(str).str.strip() == "Remediated"]) if status_col else 0
    unable_to_fix = total_lanes - offline_lanes - remediated

    summary = {
        "Activity Type": "EOD Failure",
        "# Issue Lanes": total_lanes,
        "# Offline Lanes": offline_lanes,
        "#Remediated": remediated,
        "# Unable to Fix": unable_to_fix,
    }

    # Unable to Fix breakdown (right chart)
    # Group by status (New/Offline/Pre-Existing) and Category
    unable_fix_breakdown = pd.DataFrame()
    if status_col and category_col:
        # Exclude Remediated rows
        non_remediated = df[df[status_col].astype(str).str.strip() != "Remediated"]
        pivot = pd.crosstab(
            non_remediated[status_col].astype(str).str.strip(),
            non_remediated[category_col].astype(str).str.strip()
        )
        unable_fix_breakdown = pivot

    # Repetition: count occurrences of each Host
    host_counts = df["Host"].value_counts().reset_index()
    host_counts.columns = ["Hostname", "Count"]
    host_counts["Device Type"] = host_counts["Hostname"].apply(get_device_type)
    # Get issue for each host (take first occurrence)
    if issue_sub_col:
        host_issues = df.groupby("Host")[issue_sub_col].first().reset_index()
        host_issues.columns = ["Hostname", "Issue Subcategory"]
        host_counts = host_counts.merge(host_issues, on="Hostname", how="left")
    else:
        host_counts["Issue Subcategory"] = "EOD failure"

    # Repetition chart: count of hosts by device type and repeat count
    # X-axis shows "No. of Repeats" grouped by device type, Y-axis shows count of hosts
    rep_chart_data = []
    for device in ["SCO", "REG"]:
        device_hosts = host_counts[host_counts["Device Type"] == device]
        count_groups = device_hosts.groupby("Count").size().reset_index(name="Total Count")
        for _, row in count_groups.iterrows():
            rep_chart_data.append({
                "Device Type": device,
                "No. of Repeats": int(row["Count"]),
                "Total Count": int(row["Total Count"]),
            })
    rep_chart = pd.DataFrame(rep_chart_data) if rep_chart_data else pd.DataFrame()

    # New incidents
    new_incidents = pd.DataFrame()
    if status_col:
        new_rows = df[df[status_col].astype(str).str.strip() == "New"]
        if not new_rows.empty:
            new_incidents = pd.DataFrame({
                "Hostname": new_rows["Host"].values,
                "INC": new_rows["INC"].values if "INC" in new_rows.columns else [""] * len(new_rows),
                "Issue": new_rows[issue_sub_col].values if issue_sub_col else ["EOD failure"] * len(new_rows),
            })
            new_incidents = new_incidents[new_incidents["INC"].astype(str) != "nan"]
            new_incidents = new_incidents[new_incidents["INC"].astype(str) != "NA"]
            new_incidents = new_incidents.reset_index(drop=True)

    return summary, unable_fix_breakdown, host_counts, rep_chart, new_incidents


def build_eod_ppt(summary, unable_fix_breakdown, host_counts, rep_chart, new_incidents, today_str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])

    # Header
    h = slide1.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.3))
    h.text_frame.paragraphs[0].text = f"ATC-I POS Deployment/Remediation/TA status | {today_str}"
    h.text_frame.paragraphs[0].font.size = Pt(10)
    h.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    # Title
    t = slide1.shapes.add_textbox(Inches(0.3), Inches(0.35), Inches(9.4), Inches(0.4))
    t.text_frame.paragraphs[0].text = "•  Health check Numbers – EOD"
    t.text_frame.paragraphs[0].font.size = Pt(14)
    t.text_frame.paragraphs[0].font.bold = True

    # Summary table
    cols = ["Activity Type", "# Issue Lanes", "# Offline Lanes", "#Remediated", "# Unable to Fix"]
    tbl = slide1.shapes.add_table(2, 5, Inches(0.3), Inches(0.85), Inches(9.4), Inches(0.6)).table
    for j, col in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = col
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BLUE
    for j, col in enumerate(cols):
        cell = tbl.cell(1, j)
        cell.text = str(summary[col])
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.fill.background()

    # Left chart: EOD Failure
    ct = slide1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(3.5), Inches(0.3))
    ct.text_frame.paragraphs[0].text = "EOD Failure"
    ct.text_frame.paragraphs[0].font.size = Pt(10)
    ct.text_frame.paragraphs[0].font.bold = True
    ct.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cd = CategoryChartData()
    cd.categories = ["# Issue Lanes", "# OfflineLanes", "#Remediated", "# Unable to Fix"]
    cd.add_series("EOD Failure", [summary["# Issue Lanes"], summary["# Offline Lanes"], summary["#Remediated"], summary["# Unable to Fix"]])

    cf = slide1.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.3), Inches(1.9), Inches(4.5), Inches(5.0), cd)
    chart = cf.chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = CHART_BLUE
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.number_format = '0'
    plot.data_labels.font.size = Pt(9)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    set_category_axis_horizontal(chart)

    # Right chart: Unable to Fix Status (No repeat)
    if not unable_fix_breakdown.empty:
        try:
            ct2 = slide1.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(4.0), Inches(0.3))
            ct2.text_frame.paragraphs[0].text = "Unable to Fix Status (No repeat)"
            ct2.text_frame.paragraphs[0].font.size = Pt(10)
            ct2.text_frame.paragraphs[0].font.bold = True
            ct2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            categories = unable_fix_breakdown.index.tolist()
            series_names = unable_fix_breakdown.columns.tolist()

            cd2 = CategoryChartData()
            cd2.categories = categories
            for col in series_names:
                vals = [int(unable_fix_breakdown.loc[cat, col]) for cat in categories]
                cd2.add_series(str(col), vals)

            cf2 = slide1.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.2), Inches(1.9), Inches(4.5), Inches(5.0), cd2)
            chart2 = cf2.chart
            chart2.has_legend = True
            chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart2.legend.font.size = Pt(8)
            colors = [CHART_BLUE, CHART_ORANGE, CHART_GRAY]
            for idx, s in enumerate(chart2.series):
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = colors[idx % len(colors)]
            plot2 = chart2.plots[0]
            plot2.has_data_labels = True
            plot2.data_labels.show_value = True
            plot2.data_labels.number_format = '0'
            plot2.data_labels.font.size = Pt(8)
            chart2.value_axis.has_major_gridlines = True
            chart2.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
            set_category_axis_horizontal(chart2)
        except Exception:
            pass

    # ===== SLIDE 2 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    # Header
    h2 = slide2.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.3))
    h2.text_frame.paragraphs[0].text = f"ATC-I POS Deployment/Remediation/TA status | {today_str}"
    h2.text_frame.paragraphs[0].font.size = Pt(10)
    h2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    # Repetition title
    rt = slide2.shapes.add_textbox(Inches(0.3), Inches(0.5), Inches(5.0), Inches(0.4))
    rt.text_frame.paragraphs[0].text = "•  Repetition of Lanes at EOD"
    rt.text_frame.paragraphs[0].font.size = Pt(14)
    rt.text_frame.paragraphs[0].font.bold = True

    # Repetition table (Count > 1 only)
    rep_filtered = host_counts[host_counts["Count"] > 1].reset_index(drop=True)
    if not rep_filtered.empty:
        rep_cols = ["Hostname", "Count", "Device Type", "Issue Subcategory"]
        n_rows = min(len(rep_filtered), 10)
        tbl2 = slide2.shapes.add_table(n_rows + 1, 4, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.3 + n_rows * 0.3)).table
        for j, col in enumerate(rep_cols):
            cell = tbl2.cell(0, j)
            cell.text = col
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BLUE
        for i in range(n_rows):
            row = rep_filtered.iloc[i]
            for j, col in enumerate(rep_cols):
                cell = tbl2.cell(i + 1, j)
                val = row.get(col, "")
                cell.text = str(int(val)) if col == "Count" else str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.fill.background()

    # Repetition chart (right side) — grouped by device type, x-axis = No. of Repeats
    if not rep_chart.empty:
        ct3 = slide2.shapes.add_textbox(Inches(5.5), Inches(0.4), Inches(4.0), Inches(0.3))
        ct3.text_frame.paragraphs[0].text = "Repetition Count for EOD"
        ct3.text_frame.paragraphs[0].font.size = Pt(10)
        ct3.text_frame.paragraphs[0].font.bold = True
        ct3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Build categories: "repeat_count\nDevice Type" for each bar
        categories = []
        values = []
        for _, row in rep_chart.iterrows():
            categories.append(f"{int(row['No. of Repeats'])}\n{row['Device Type']}")
            values.append(int(row["Total Count"]))

        cd3 = CategoryChartData()
        cd3.categories = categories
        cd3.add_series("Total Count", values)

        cf3 = slide2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.2), Inches(0.7), Inches(4.5), Inches(3.0), cd3)
        chart3 = cf3.chart
        chart3.has_legend = False
        chart3.series[0].format.fill.solid()
        chart3.series[0].format.fill.fore_color.rgb = CHART_BLUE
        plot3 = chart3.plots[0]
        plot3.has_data_labels = True
        plot3.data_labels.show_value = True
        plot3.data_labels.number_format = '0'
        plot3.data_labels.font.size = Pt(9)
        chart3.value_axis.has_major_gridlines = True
        chart3.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
        set_category_axis_horizontal(chart3)

    # New incidents title
    new_top = Inches(1.0 + (min(len(rep_filtered), 10) + 1) * 0.3 + 0.5) if not rep_filtered.empty else Inches(3.5)
    nt = slide2.shapes.add_textbox(Inches(0.3), new_top, Inches(5.0), Inches(0.4))
    nt.text_frame.paragraphs[0].text = "•  New incidents raised"
    nt.text_frame.paragraphs[0].font.size = Pt(14)
    nt.text_frame.paragraphs[0].font.bold = True

    # New incidents table
    if not new_incidents.empty:
        n_new = min(len(new_incidents), 10)
        new_tbl_top = new_top + Inches(0.5)
        tbl3 = slide2.shapes.add_table(n_new + 1, 3, Inches(0.3), new_tbl_top, Inches(4.5), Inches(0.3 + n_new * 0.3)).table
        for j, col in enumerate(["Hostname", "INC", "Issue"]):
            cell = tbl3.cell(0, j)
            cell.text = col
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BLUE
        for i in range(n_new):
            row = new_incidents.iloc[i]
            for j, col in enumerate(["Hostname", "INC", "Issue"]):
                cell = tbl3.cell(i + 1, j)
                cell.text = str(row.get(col, ""))
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.fill.background()

    return prs


