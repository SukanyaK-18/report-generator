"""
POS Readiness Report Generator - Streamlit Web App
Generates a 2-slide PowerPoint for POS-REG or POS-SCO Weekly Health Check.
"""

import io
import pandas as pd
import numpy as np
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from datetime import datetime

HEADER_BLUE = RGBColor(0x44, 0x72, 0xC4)
CHART_BLUE = RGBColor(0x44, 0x72, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def get_device_type(hostname):
    hostname = str(hostname).upper()
    if "SCO" in hostname:
        return "SCO"
    elif "REG" in hostname:
        return "REG"
    return "Unknown"


def set_category_axis_horizontal(chart):
    catAx = chart._element.find(".//" + qn("c:catAx"))
    if catAx is not None:
        existing = catAx.find(qn("c:txPr"))
        if existing is not None:
            catAx.remove(existing)
        txPr = etree.SubElement(catAx, qn("c:txPr"))
        bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
        bodyPr.set("rot", "0")
        bodyPr.set("vert", "horz")
        bodyPr.set("spcFirstLastPara", "1")
        bodyPr.set("wrap", "square")
        etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        defRPr.set("sz", "700")


def compute_pos_data(df, report_type):
    """Compute metrics for POS-REG or POS-SCO readiness."""
    total_lanes = len(df)
    offline = len(df[df["Status"].astype(str).str.strip() == "Offline"])
    false_positive = len(df[df["Status"].astype(str).str.strip() == "False Positive"])
    resolved = len(df[df["Status"].astype(str).str.strip() == "Completed"])
    # Follow-up needed = In Progress with Follow-up = Yes
    in_progress = df[df["Status"].astype(str).str.strip().str.lower().str.startswith("in progress")]
    follow_up = len(in_progress[in_progress["Follow-up needed"].astype(str).str.strip() == "Yes"])

    summary = {
        "Activity Type": report_type,
        "# Issue Lanes": total_lanes,
        "#Offline": offline,
        "# False Positive": false_positive,
        "# Resolved": resolved,
        "# Follow-up Needed": follow_up,
    }

    # Unable to fix breakdown: In Progress rows (excluding Offline category) grouped by Issue Sub-Category + Category
    unable_to_fix = in_progress.copy()
    unable_fix_chart = pd.DataFrame()
    if not unable_to_fix.empty and "Issue Sub-Category" in unable_to_fix.columns:
        # Exclude Offline category AND Offline sub-category from "Unable to fix" chart
        unable_to_fix = unable_to_fix[unable_to_fix["Category"].astype(str).str.strip().str.lower() != "offline"]
        unable_to_fix = unable_to_fix[unable_to_fix["Issue Sub-Category"].astype(str).str.strip().str.lower() != "offline"]
        # Also exclude Splunk Issue (false positives already resolved)
        unable_to_fix = unable_to_fix[unable_to_fix["Category"].astype(str).str.strip().str.lower() != "splunk issue"]
        if not unable_to_fix.empty:
            # Group by Issue Sub-Category + Category, sort by Category to group visually
            grouped = unable_to_fix.groupby(["Category", "Issue Sub-Category"]).size().reset_index(name="Count")
            grouped = grouped.sort_values(["Category", "Count"], ascending=[True, False])
            unable_fix_chart = grouped

    # Repetition: count of "unable to fix" (In Progress) hosts
    host_counts = pd.DataFrame()
    if not in_progress.empty:
        counts = in_progress["Host"].value_counts().reset_index()
        counts.columns = ["Hostname", "Count"]
        counts["Device Type"] = counts["Hostname"].apply(get_device_type)
        # Get issue category for each host
        if "Category" in in_progress.columns:
            host_cat = in_progress.groupby("Host")["Issue Sub-Category"].first().reset_index()
            host_cat.columns = ["Hostname", "Issue Category"]
            counts = counts.merge(host_cat, on="Hostname", how="left")
        else:
            counts["Issue Category"] = ""
        host_counts = counts

    # Repetition chart data
    rep_chart_data = []
    if not host_counts.empty:
        for device in host_counts["Device Type"].unique():
            device_hosts = host_counts[host_counts["Device Type"] == device]
            count_groups = device_hosts.groupby("Count").size().reset_index(name="Total Count")
            for _, row in count_groups.iterrows():
                rep_chart_data.append({
                    "Device Type": device,
                    "No. of Repeats": int(row["Count"]),
                    "Total Count": int(row["Total Count"]),
                })
    rep_chart = pd.DataFrame(rep_chart_data) if rep_chart_data else pd.DataFrame()

    return summary, unable_fix_chart, host_counts, rep_chart


def build_pos_ppt(summary, unable_fix_chart, host_counts, rep_chart, report_type, today_str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1 =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])

    # Header
    h = slide1.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.3))
    h.text_frame.paragraphs[0].text = f"ATC-I POS Deployment/Remediation/TA status | {today_str}"
    h.text_frame.paragraphs[0].font.size = Pt(10)
    h.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    # Title
    t = slide1.shapes.add_textbox(Inches(0.3), Inches(0.35), Inches(9.4), Inches(0.4))
    t.text_frame.paragraphs[0].text = f"•  Health check Numbers – POS - {report_type}"
    t.text_frame.paragraphs[0].font.size = Pt(14)
    t.text_frame.paragraphs[0].font.bold = True

    # Summary table
    cols = ["Activity Type", "# Issue Lanes", "#Offline", "# False Positive", "# Resolved", "# Follow-up Needed"]
    tbl = slide1.shapes.add_table(2, 6, Inches(0.2), Inches(0.85), Inches(9.6), Inches(0.6)).table
    for j, col in enumerate(cols):
        cell = tbl.cell(0, j)
        cell.text = col
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BLUE
    for j, col in enumerate(cols):
        cell = tbl.cell(1, j)
        cell.text = str(summary[col])
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.fill.background()

    # Left chart
    ct = slide1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(3.5), Inches(0.3))
    ct.text_frame.paragraphs[0].text = report_type
    ct.text_frame.paragraphs[0].font.size = Pt(10)
    ct.text_frame.paragraphs[0].font.bold = True
    ct.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    cd = CategoryChartData()
    cd.categories = ["# Issue Lanes", "#Offline", "# False Positive", "# Resolved", "# Follow-up\nNeeded"]
    cd.add_series(report_type, [
        summary["# Issue Lanes"], summary["#Offline"], summary["# False Positive"],
        summary["# Resolved"], summary["# Follow-up Needed"]
    ])

    cf = slide1.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.3), Inches(1.9), Inches(4.5), Inches(5.0), cd)
    chart = cf.chart
    chart.has_legend = False
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = CHART_BLUE
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.number_format = '0'
    plot.data_labels.font.size = Pt(9)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    set_category_axis_horizontal(chart)

    # Right chart: Unable to fix
    if not unable_fix_chart.empty:
        try:
            ct2 = slide1.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(4.0), Inches(0.3))
            ct2.text_frame.paragraphs[0].text = "Unable to fix"
            ct2.text_frame.paragraphs[0].font.size = Pt(10)
            ct2.text_frame.paragraphs[0].font.bold = True
            ct2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Categories = "Issue Sub-Category\nCategory"
            categories = [f"{row['Issue Sub-Category']}\n{row['Category']}" for _, row in unable_fix_chart.iterrows()]
            values = [int(row["Count"]) for _, row in unable_fix_chart.iterrows()]

            cd2 = CategoryChartData()
            cd2.categories = categories
            cd2.add_series("Count", values)

            cf2 = slide1.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5.2), Inches(1.9), Inches(4.5), Inches(5.0), cd2)
            chart2 = cf2.chart
            chart2.has_legend = False
            chart2.series[0].format.fill.solid()
            chart2.series[0].format.fill.fore_color.rgb = CHART_BLUE
            plot2 = chart2.plots[0]
            plot2.has_data_labels = True
            plot2.data_labels.show_value = True
            plot2.data_labels.number_format = '0'
            plot2.data_labels.font.size = Pt(8)
            chart2.value_axis.has_major_gridlines = True
            chart2.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
            set_category_axis_horizontal(chart2)

            # Add "In Progress" label below the chart as x-axis title
            ip_label = slide1.shapes.add_textbox(Inches(6.5), Inches(6.7), Inches(2.0), Inches(0.3))
            ip_label.text_frame.paragraphs[0].text = "In Progress"
            ip_label.text_frame.paragraphs[0].font.size = Pt(9)
            ip_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        except Exception:
            pass

    # ===== SLIDE 2 =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])

    # Header
    h2 = slide2.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.3))
    h2.text_frame.paragraphs[0].text = f"ATC-I POS Deployment/Remediation/TA status | {today_str}"
    h2.text_frame.paragraphs[0].font.size = Pt(10)
    h2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x59, 0x59, 0x59)

    # Repetition title
    rt = slide2.shapes.add_textbox(Inches(0.3), Inches(0.4), Inches(9.0), Inches(0.4))
    rt.text_frame.paragraphs[0].text = f"•  Repetition of Unable to fix Lanes – {report_type}"
    rt.text_frame.paragraphs[0].font.size = Pt(14)
    rt.text_frame.paragraphs[0].font.bold = True

    # Repetition table (Count > 1 only)
    rep_filtered = host_counts[host_counts["Count"] > 1].reset_index(drop=True) if not host_counts.empty else pd.DataFrame()
    if not rep_filtered.empty:
        rep_cols = ["Hostname", "Count", "Device Type", "Issue Category"]
        n_rows = min(len(rep_filtered), 10)
        tbl2 = slide2.shapes.add_table(n_rows + 1, 4, Inches(0.3), Inches(0.9), Inches(5.5), Inches(0.3 + n_rows * 0.3)).table
        for j, col in enumerate(rep_cols):
            cell = tbl2.cell(0, j)
            cell.text = col
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BLUE
        for i in range(n_rows):
            row = rep_filtered.iloc[i]
            for j, col in enumerate(rep_cols):
                cell = tbl2.cell(i + 1, j)
                val = row.get(col, "")
                cell.text = str(int(val)) if col == "Count" else str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.fill.background()

    # Repetition chart
    if not rep_chart.empty:
        ct3 = slide2.shapes.add_textbox(Inches(2.0), Inches(1.0 + (min(len(rep_filtered), 10) + 1) * 0.3 + 0.3) if not rep_filtered.empty else Inches(1.5), Inches(6.0), Inches(0.3))
        ct3.text_frame.paragraphs[0].text = "Repetition of Unable to fix Lanes"
        ct3.text_frame.paragraphs[0].font.size = Pt(10)
        ct3.text_frame.paragraphs[0].font.bold = True
        ct3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        categories = [f"{int(row['No. of Repeats'])}\n{row['Device Type']}" for _, row in rep_chart.iterrows()]
        values = [int(row["Total Count"]) for _, row in rep_chart.iterrows()]

        cd3 = CategoryChartData()
        cd3.categories = categories
        cd3.add_series("Total Count", values)

        chart_top = Inches(1.0 + (min(len(rep_filtered), 10) + 1) * 0.3 + 0.6) if not rep_filtered.empty else Inches(1.8)
        cf3 = slide2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), chart_top, Inches(7.0), Inches(4.0), cd3)
        chart3 = cf3.chart
        chart3.has_legend = False
        chart3.series[0].format.fill.solid()
        chart3.series[0].format.fill.fore_color.rgb = CHART_BLUE
        plot3 = chart3.plots[0]
        plot3.has_data_labels = True
        plot3.data_labels.show_value = True
        plot3.data_labels.number_format = '0'
        plot3.data_labels.font.size = Pt(9)
        chart3.value_axis.has_major_gridlines = True
        chart3.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
        set_category_axis_horizontal(chart3)

    return prs


